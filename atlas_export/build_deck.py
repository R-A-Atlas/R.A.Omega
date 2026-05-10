"""Build PowerPoint from POST /query-shaped analysis JSON."""

from __future__ import annotations

import json
import re
from datetime import date
from pathlib import Path
from typing import Any

ROOT = Path(__file__).resolve().parent.parent
DECKS_DIR = ROOT / "atlas_vault" / "03-Outputs" / "Decks"


def _tick(envelope: dict[str, Any]) -> str:
    pq = envelope.get("parsed_query") or {}
    if isinstance(pq, dict):
        t = pq.get("tickers")
        if isinstance(t, (list, tuple)) and t:
            return str(t[0]).upper().strip()[:12]
    fr = envelope.get("final_report") or {}
    if isinstance(fr, dict) and fr.get("ticker"):
        return str(fr["ticker"]).upper().strip()[:12]
    return "ATLAS"


def _txt(x: Any, limit: int = 6000) -> str:
    if x is None:
        return ""
    if isinstance(x, (dict, list)):
        try:
            s = json.dumps(x, indent=2)[:limit]
        except (TypeError, ValueError):
            s = str(x)[:limit]
    else:
        s = str(x).strip()
    return s[:limit]


def _trade_plan_lines(fr: dict[str, Any]) -> list[str]:
    tp = fr.get("trade_plan")
    if isinstance(tp, str) and tp.strip():
        return [tp.strip()]
    if isinstance(tp, dict):
        lines = []
        for k, v in tp.items():
            if v not in (None, ""):
                lines.append(f"{k}: {v}")
        return lines[:24]
    return []


def _scenario_lines(env: dict[str, Any]) -> list[str]:
    arr = env.get("scenarios")
    if not isinstance(arr, list):
        return []
    out = []
    for sc in arr[:6]:
        if not isinstance(sc, dict):
            out.append(str(sc))
            continue
        lab = sc.get("label") or sc.get("name") or ""
        prob = sc.get("probability") or sc.get("prob")
        trig = sc.get("trigger") or ""
        oc = sc.get("outcome") or ""
        line = " · ".join(str(p) for p in (lab, prob, trig, oc) if p not in (None, ""))
        if line:
            out.append(line)
    return out


def _risk_lines(fr: dict[str, Any]) -> list[str]:
    kr = fr.get("key_risks")
    if isinstance(kr, list):
        return [str(x.get("risk") or x.get("text") or x)[:400] for x in kr[:12] if x]
    if isinstance(kr, str) and kr.strip():
        return [kr.strip()[:2000]]
    rt = fr.get("risks_and_tripwires")
    if isinstance(rt, list):
        return [str(x.get("risk") or x)[:400] for x in rt[:12]]
    return []


def _catalyst_lines(fr: dict[str, Any]) -> list[str]:
    c = fr.get("catalysts_timeline")
    if isinstance(c, list):
        rows = []
        for x in c[:16]:
            if isinstance(x, dict):
                d = x.get("date") or x.get("when") or ""
                e = x.get("event") or x.get("label") or ""
                rows.append(f"{d} — {e}".strip(" —"))
            elif isinstance(x, str) and x.strip():
                rows.append(x.strip())
        return rows
    if isinstance(c, str) and c.strip():
        return [ln.strip() for ln in c.splitlines() if ln.strip()][:20]
    return []


def write_query_envelope_pptx(
    envelope: dict[str, Any],
    dest: Path | None = None,
) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    DECKS_DIR.mkdir(parents=True, exist_ok=True)
    tk = re.sub(r"[^\w\-]+", "_", _tick(envelope))[:32] or "ATLAS"
    target = dest or (DECKS_DIR / f"{tk}_{date.today().strftime('%Y-%m-%d')}.pptx")

    fr = envelope.get("final_report") if isinstance(envelope.get("final_report"), dict) else {}
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    def add_title_slide(title: str, subtitle: str = "") -> None:
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

    def add_bullets(title: str, bullets: list[str]) -> None:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title[:200]
        body = slide.shapes.placeholders[1].text_frame
        lines = bullets[:24] if bullets else ["—"]
        body.text = lines[0][:500]
        for line in lines[1:]:
            p = body.add_paragraph()
            p.text = line[:500]
            p.level = 0
            p.font.size = Pt(14)

    tkr = _tick(envelope)
    q = _txt(envelope.get("query"), 200)
    add_title_slide(f"ATLAS — {tkr}", q or "Equity / finance analysis")

    tldr = _txt(envelope.get("tldr") or fr.get("tldr"), 3500)
    add_bullets("TLDR", [tldr] if tldr else ["—"])

    exe = _txt(fr.get("executive_summary"), 4000)
    add_bullets("Executive summary", [exe] if exe else ["—"])

    tp_lines = _trade_plan_lines(fr)
    add_bullets("Trade plan", tp_lines if tp_lines else [_txt(fr.get("trade_plan"), 800)])

    scen = _scenario_lines(envelope)
    half = max(1, len(scen) // 2) if scen else 1
    add_bullets("Scenarios (part 1)", scen[:half] if scen else ["—"])
    if len(scen) > half:
        add_bullets("Scenarios (part 2)", scen[half:])

    pl = fr.get("price_levels")
    plines: list[str] = []
    if isinstance(pl, dict):
        for k, v in list(pl.items())[:20]:
            plines.append(f"{k}: {v}")
    elif isinstance(pl, str):
        plines = [ln.strip() for ln in pl.splitlines() if ln.strip()][:20]
    add_bullets("Price levels", plines if plines else ["—"])

    risks = _risk_lines(fr)
    add_bullets("Risk factors", risks if risks else ["See full report JSON"])

    cats = _catalyst_lines(fr)
    add_bullets("Catalyst timeline", cats if cats else ["—"])

    add_bullets(
        "Disclaimer",
        [
            "Informational only. Not investment advice.",
            "Verify data independently before acting.",
        ],
    )

    prs.save(str(target.resolve()))
    return target


__all__ = ["write_query_envelope_pptx", "DECKS_DIR"]
